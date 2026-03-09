from __future__ import annotations

import math
import os
import io
import re
import base64
import time
import warnings
warnings.filterwarnings("ignore")

from flask import Flask, render_template, request, jsonify, send_file
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from xml.etree import ElementTree as ET

app = Flask(__name__)

GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

SLIDE_W = 13.33
SLIDE_H = 7.5
EMU_PER_IN = 914400


# ── SVG helpers ───────────────────────────────────────────────────────────────

def _dedup_attrs(tag_match):
    """Remove duplicate attributes within a single SVG tag."""
    tag_str = tag_match.group(0)
    # Find all attr="value" pairs
    attrs = re.findall(r'''([\w\-:]+)\s*=\s*("[^"]*"|'[^']*')''', tag_str)
    if not attrs:
        return tag_str
    seen = {}
    for name, val in attrs:
        seen[name] = val  # last one wins
    # Check if there were duplicates
    if len(seen) == len(attrs):
        return tag_str  # no duplicates
    # Rebuild the tag
    tag_name_m = re.match(r'<([\w\-:]+)', tag_str)
    if not tag_name_m:
        return tag_str
    tag_name = tag_name_m.group(1)
    self_close = tag_str.rstrip().endswith("/>")
    parts = [f"<{tag_name}"]
    for name, val in seen.items():
        parts.append(f' {name}={val}')
    parts.append("/>" if self_close else ">")
    return "".join(parts)


def extract_svg(raw: str) -> str | None:
    if not raw:
        return None
    s = re.sub(r"```(?:svg|xml|html)?", "", raw).replace("```", "").strip()
    a = s.lower().find("<svg")
    b = s.lower().rfind("</svg>")
    if a != -1 and b != -1 and b > a:
        svg = s[a : b + 6]
        svg = re.sub(r"&(?!(?:amp|lt|gt|quot|apos|#\d+|#x[0-9a-fA-F]+);)", "&amp;", svg)
        if 'xmlns=' not in svg.split('>')[0]:
            svg = svg.replace('<svg', '<svg xmlns="http://www.w3.org/2000/svg"', 1)
        # Remove duplicate attributes (Gemini sometimes emits them)
        svg = re.sub(r'<[^>]+>', _dedup_attrs, svg)
        return svg
    return None


def hex_to_rgb(h: str) -> tuple:
    h = h.lstrip("#")
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


# ── SVG attr parsing helpers ─────────────────────────────────────────────────

_SVG_COLORS = {
    "black": "#000000", "white": "#FFFFFF", "red": "#FF0000",
    "green": "#008000", "blue": "#0000FF", "yellow": "#FFFF00",
    "gray": "#808080", "grey": "#808080", "silver": "#C0C0C0",
    "navy": "#000080", "teal": "#008080", "orange": "#FFA500",
    "purple": "#800080", "maroon": "#800000", "aqua": "#00FFFF",
    "lime": "#00FF00", "olive": "#808000", "fuchsia": "#FF00FF",
    "darkgray": "#A9A9A9", "darkgrey": "#A9A9A9",
    "lightgray": "#D3D3D3", "lightgrey": "#D3D3D3",
    "dimgray": "#696969", "dimgrey": "#696969",
    "slategray": "#708090", "slategrey": "#708090",
    "darkslategray": "#2F4F4F", "darkslategrey": "#2F4F4F",
    "gainsboro": "#DCDCDC", "whitesmoke": "#F5F5F5",
    "cornflowerblue": "#6495ED", "steelblue": "#4682B4",
    "dodgerblue": "#1E90FF", "royalblue": "#4169E1",
    "midnightblue": "#191970", "darkblue": "#00008B",
    "indigo": "#4B0082", "coral": "#FF7F50", "tomato": "#FF6347",
    "salmon": "#FA8072", "gold": "#FFD700", "khaki": "#F0E68C",
}


def _ga(elem, attr, default=""):
    """Get SVG attribute, checking inline style as fallback."""
    v = elem.get(attr, "")
    if v:
        return v
    for part in elem.get("style", "").split(";"):
        if ":" in part:
            k, val = part.split(":", 1)
            if k.strip() == attr:
                return val.strip()
    return default


def _color(val):
    """Parse SVG color string -> RGBColor or None."""
    if not val or val in ("none", "transparent", "inherit", "currentColor"):
        return None
    val = val.strip()
    # Named color
    if val.lower() in _SVG_COLORS:
        val = _SVG_COLORS[val.lower()]
    # rgb(r,g,b)
    m = re.match(r"rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)", val)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    # Hex
    try:
        return RGBColor(*hex_to_rgb(val))
    except Exception:
        return None


def _num(val, default=0.0):
    """Parse numeric, stripping px/em/etc."""
    if not val:
        return default
    try:
        return float(re.sub(r"[a-z%]+", "", str(val), flags=re.I).strip())
    except Exception:
        return default


def _tag(elem):
    t = elem.tag
    return t.split("}")[1] if "}" in t else t


def _parse_transform(transform):
    """Parse translate() from a transform string."""
    if not transform:
        return 0.0, 0.0
    m = re.search(r"translate\(\s*([-\d.e]+)[,\s]+([-\d.e]+)", transform)
    if m:
        return float(m.group(1)), float(m.group(2))
    m = re.search(r"translate\(\s*([-\d.e]+)\s*\)", transform)
    if m:
        return float(m.group(1)), 0.0
    return 0.0, 0.0


def _opacity(elem):
    """Get opacity value (0.0 to 1.0). Returns 1.0 if not set."""
    v = _ga(elem, "opacity", "1")
    try:
        return float(v)
    except Exception:
        return 1.0


# ── Bezier math helpers ──────────────────────────────────────────────────────

def _cubic_bezier_point(t, p0, p1, p2, p3):
    """Evaluate cubic bezier at parameter t."""
    u = 1 - t
    return (u**3 * p0[0] + 3*u**2*t * p1[0] + 3*u*t**2 * p2[0] + t**3 * p3[0],
            u**3 * p0[1] + 3*u**2*t * p1[1] + 3*u*t**2 * p2[1] + t**3 * p3[1])


def _quad_bezier_point(t, p0, p1, p2):
    """Evaluate quadratic bezier at parameter t."""
    u = 1 - t
    return (u**2 * p0[0] + 2*u*t * p1[0] + t**2 * p2[0],
            u**2 * p0[1] + 2*u*t * p1[1] + t**2 * p2[1])


def _arc_to_points(cx, cy, rx, ry, start_angle, sweep_angle, steps=12):
    """Sample points along an SVG arc."""
    pts = []
    for i in range(steps + 1):
        angle = start_angle + sweep_angle * i / steps
        pts.append((cx + rx * math.cos(angle), cy + ry * math.sin(angle)))
    return pts


def _svg_arc_to_center(x1, y1, rx, ry, phi, fa, fs, x2, y2):
    """Convert SVG arc endpoint parameterization to center parameterization."""
    cos_phi = math.cos(math.radians(phi))
    sin_phi = math.sin(math.radians(phi))

    dx2 = (x1 - x2) / 2
    dy2 = (y1 - y2) / 2
    x1p = cos_phi * dx2 + sin_phi * dy2
    y1p = -sin_phi * dx2 + cos_phi * dy2

    rx = abs(rx)
    ry = abs(ry)
    if rx == 0 or ry == 0:
        return [(x2, y2)]

    lam = (x1p**2) / (rx**2) + (y1p**2) / (ry**2)
    if lam > 1:
        s = math.sqrt(lam)
        rx *= s
        ry *= s

    num = max(0, rx**2 * ry**2 - rx**2 * y1p**2 - ry**2 * x1p**2)
    den = rx**2 * y1p**2 + ry**2 * x1p**2
    if den == 0:
        return [(x2, y2)]
    sq = math.sqrt(num / den)
    if fa == fs:
        sq = -sq

    cxp = sq * rx * y1p / ry
    cyp = -sq * ry * x1p / rx

    cx_c = cos_phi * cxp - sin_phi * cyp + (x1 + x2) / 2
    cy_c = sin_phi * cxp + cos_phi * cyp + (y1 + y2) / 2

    theta1 = math.atan2((y1p - cyp) / ry, (x1p - cxp) / rx)
    dtheta = math.atan2((-y1p - cyp) / ry, (-x1p - cxp) / rx) - theta1

    if fs == 0 and dtheta > 0:
        dtheta -= 2 * math.pi
    elif fs == 1 and dtheta < 0:
        dtheta += 2 * math.pi

    return _arc_to_points(cx_c, cy_c, rx, ry, theta1, dtheta, steps=16)


# ── Full SVG path parser (with curves) ───────────────────────────────────────

def _parse_path_full(d, ox, oy):
    """Parse SVG path d attribute → list of (x, y) points.
    Supports M/L/H/V/Z/C/S/Q/T/A (absolute and relative).
    Curves are approximated as line segments."""
    result = []
    cur_x, cur_y = 0.0, 0.0
    start_x, start_y = 0.0, 0.0
    last_cp_x, last_cp_y = 0.0, 0.0  # for S/T smooth continuations
    last_cmd = ""

    tokens = re.findall(
        r"[MmLlHhVvZzCcSsQqTtAa]|[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?",
        d
    )
    i = 0
    cmd = "M"

    def _float(idx):
        return float(tokens[idx])

    while i < len(tokens):
        t = tokens[i]
        if t.isalpha():
            cmd = t
            i += 1
            if cmd in "Zz":
                result.append((start_x + ox, start_y + oy))
                cur_x, cur_y = start_x, start_y
                last_cmd = cmd
                continue
        else:
            # Implicit repeat of previous command
            pass

        try:
            if cmd == "M":
                cur_x, cur_y = _float(i), _float(i+1)
                start_x, start_y = cur_x, cur_y
                result.append((cur_x + ox, cur_y + oy))
                i += 2
                cmd = "L"  # subsequent coords are lines
            elif cmd == "m":
                cur_x += _float(i); cur_y += _float(i+1)
                start_x, start_y = cur_x, cur_y
                result.append((cur_x + ox, cur_y + oy))
                i += 2
                cmd = "l"
            elif cmd == "L":
                cur_x, cur_y = _float(i), _float(i+1)
                result.append((cur_x + ox, cur_y + oy))
                i += 2
            elif cmd == "l":
                cur_x += _float(i); cur_y += _float(i+1)
                result.append((cur_x + ox, cur_y + oy))
                i += 2
            elif cmd == "H":
                cur_x = _float(i)
                result.append((cur_x + ox, cur_y + oy))
                i += 1
            elif cmd == "h":
                cur_x += _float(i)
                result.append((cur_x + ox, cur_y + oy))
                i += 1
            elif cmd == "V":
                cur_y = _float(i)
                result.append((cur_x + ox, cur_y + oy))
                i += 1
            elif cmd == "v":
                cur_y += _float(i)
                result.append((cur_x + ox, cur_y + oy))
                i += 1
            elif cmd == "C":
                x1, y1 = _float(i), _float(i+1)
                x2, y2 = _float(i+2), _float(i+3)
                ex, ey = _float(i+4), _float(i+5)
                for t_val in [0.25, 0.5, 0.75, 1.0]:
                    px, py = _cubic_bezier_point(
                        t_val, (cur_x, cur_y), (x1, y1), (x2, y2), (ex, ey))
                    result.append((px + ox, py + oy))
                last_cp_x, last_cp_y = x2, y2
                cur_x, cur_y = ex, ey
                i += 6
            elif cmd == "c":
                x1, y1 = cur_x + _float(i), cur_y + _float(i+1)
                x2, y2 = cur_x + _float(i+2), cur_y + _float(i+3)
                ex, ey = cur_x + _float(i+4), cur_y + _float(i+5)
                for t_val in [0.25, 0.5, 0.75, 1.0]:
                    px, py = _cubic_bezier_point(
                        t_val, (cur_x, cur_y), (x1, y1), (x2, y2), (ex, ey))
                    result.append((px + ox, py + oy))
                last_cp_x, last_cp_y = x2, y2
                cur_x, cur_y = ex, ey
                i += 6
            elif cmd == "S":
                # Smooth cubic: reflect last control point
                x1 = 2 * cur_x - last_cp_x if last_cmd in "CcSs" else cur_x
                y1 = 2 * cur_y - last_cp_y if last_cmd in "CcSs" else cur_y
                x2, y2 = _float(i), _float(i+1)
                ex, ey = _float(i+2), _float(i+3)
                for t_val in [0.25, 0.5, 0.75, 1.0]:
                    px, py = _cubic_bezier_point(
                        t_val, (cur_x, cur_y), (x1, y1), (x2, y2), (ex, ey))
                    result.append((px + ox, py + oy))
                last_cp_x, last_cp_y = x2, y2
                cur_x, cur_y = ex, ey
                i += 4
            elif cmd == "s":
                x1 = 2 * cur_x - last_cp_x if last_cmd in "CcSs" else cur_x
                y1 = 2 * cur_y - last_cp_y if last_cmd in "CcSs" else cur_y
                x2 = cur_x + _float(i); y2 = cur_y + _float(i+1)
                ex = cur_x + _float(i+2); ey = cur_y + _float(i+3)
                for t_val in [0.25, 0.5, 0.75, 1.0]:
                    px, py = _cubic_bezier_point(
                        t_val, (cur_x, cur_y), (x1, y1), (x2, y2), (ex, ey))
                    result.append((px + ox, py + oy))
                last_cp_x, last_cp_y = x2, y2
                cur_x, cur_y = ex, ey
                i += 4
            elif cmd == "Q":
                x1, y1 = _float(i), _float(i+1)
                ex, ey = _float(i+2), _float(i+3)
                for t_val in [0.25, 0.5, 0.75, 1.0]:
                    px, py = _quad_bezier_point(
                        t_val, (cur_x, cur_y), (x1, y1), (ex, ey))
                    result.append((px + ox, py + oy))
                last_cp_x, last_cp_y = x1, y1
                cur_x, cur_y = ex, ey
                i += 4
            elif cmd == "q":
                x1 = cur_x + _float(i); y1 = cur_y + _float(i+1)
                ex = cur_x + _float(i+2); ey = cur_y + _float(i+3)
                for t_val in [0.25, 0.5, 0.75, 1.0]:
                    px, py = _quad_bezier_point(
                        t_val, (cur_x, cur_y), (x1, y1), (ex, ey))
                    result.append((px + ox, py + oy))
                last_cp_x, last_cp_y = x1, y1
                cur_x, cur_y = ex, ey
                i += 4
            elif cmd == "T":
                x1 = 2 * cur_x - last_cp_x if last_cmd in "QqTt" else cur_x
                y1 = 2 * cur_y - last_cp_y if last_cmd in "QqTt" else cur_y
                ex, ey = _float(i), _float(i+1)
                for t_val in [0.25, 0.5, 0.75, 1.0]:
                    px, py = _quad_bezier_point(
                        t_val, (cur_x, cur_y), (x1, y1), (ex, ey))
                    result.append((px + ox, py + oy))
                last_cp_x, last_cp_y = x1, y1
                cur_x, cur_y = ex, ey
                i += 2
            elif cmd == "t":
                x1 = 2 * cur_x - last_cp_x if last_cmd in "QqTt" else cur_x
                y1 = 2 * cur_y - last_cp_y if last_cmd in "QqTt" else cur_y
                ex = cur_x + _float(i); ey = cur_y + _float(i+1)
                for t_val in [0.25, 0.5, 0.75, 1.0]:
                    px, py = _quad_bezier_point(
                        t_val, (cur_x, cur_y), (x1, y1), (ex, ey))
                    result.append((px + ox, py + oy))
                last_cp_x, last_cp_y = x1, y1
                cur_x, cur_y = ex, ey
                i += 2
            elif cmd == "A":
                arx, ary = _float(i), _float(i+1)
                rotation = _float(i+2)
                fa = int(_float(i+3))
                fs = int(_float(i+4))
                ex, ey = _float(i+5), _float(i+6)
                arc_pts = _svg_arc_to_center(cur_x, cur_y, arx, ary,
                                             rotation, fa, fs, ex, ey)
                for px, py in arc_pts[1:]:  # skip first (= current)
                    result.append((px + ox, py + oy))
                cur_x, cur_y = ex, ey
                i += 7
            elif cmd == "a":
                arx, ary = _float(i), _float(i+1)
                rotation = _float(i+2)
                fa = int(_float(i+3))
                fs = int(_float(i+4))
                ex = cur_x + _float(i+5); ey = cur_y + _float(i+6)
                arc_pts = _svg_arc_to_center(cur_x, cur_y, arx, ary,
                                             rotation, fa, fs, ex, ey)
                for px, py in arc_pts[1:]:
                    result.append((px + ox, py + oy))
                cur_x, cur_y = ex, ey
                i += 7
            else:
                i += 1
        except (IndexError, ValueError):
            i += 1

        last_cmd = cmd

    return result


# ── SVG → Editable PPTX converter ────────────────────────────────────────────

def _fix_svg_xml(svg_str):
    """Fix common SVG XML issues that cause parsing errors."""
    s = svg_str
    # Strip namespace declarations
    s = re.sub(r"\sxmlns[^\"]*\"[^\"]*\"", "", s)
    # Remove duplicate attributes
    s = re.sub(r'<[^>]+>', _dedup_attrs, s)
    # SVG void elements (should be self-closing, never have children)
    _void = ("line", "rect", "circle", "ellipse", "polygon", "polyline",
             "path", "use", "image", "br", "hr", "img")
    for tag in _void:
        # Fix <tag ...> (not self-closed) that has no matching </tag>
        # by adding self-close slash
        pattern = rf'<({tag}\b[^>]*[^/])>'
        def _fix_void(m):
            inner = m.group(1)
            return f'<{inner}/>'
        s = re.sub(pattern, _fix_void, s)
        # Remove any erroneous closing tags for void elements
        s = re.sub(rf'</{tag}\s*>', '', s)
    # Remove CDATA
    s = re.sub(r'<!\[CDATA\[.*?\]\]>', '', s, flags=re.DOTALL)
    return s


def build_editable_pptx(svg_str, bg_color, text_color):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Default background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(*hex_to_rgb(bg_color))

    # Parse SVG with error recovery
    clean = _fix_svg_xml(svg_str)
    try:
        root = ET.fromstring(clean)
    except ET.ParseError:
        # Aggressive fix: wrap in try with html.parser-style recovery
        # Strip everything after last </svg> and before first <svg
        m = re.search(r'(<svg\b.*</svg>)', clean, re.DOTALL)
        if m:
            clean = m.group(1)
        # Try removing problematic tags entirely
        clean = re.sub(r'<(?!svg|/svg|rect|/rect|text|/text|tspan|/tspan|line|circle|ellipse|polygon|g|/g)[^>]*/?>', '', clean)
        root = ET.fromstring(clean)

    vb = root.get("viewBox", "0 0 1920 1080").split()
    vb_w = float(vb[2]) if len(vb) > 2 else 1920
    vb_h = float(vb[3]) if len(vb) > 3 else 1080
    sx = SLIDE_W / vb_w   # inches per SVG unit
    sy = SLIDE_H / vb_h

    ctx = {
        "sx": sx, "sy": sy, "slide": slide,
        "def_tc": text_color, "bg_done": False,
        "vb_w": vb_w, "vb_h": vb_h,
        # Track rects so text can be matched to containing boxes
        "rects": [],
    }

    # First pass: collect all rects
    _collect_rects(root, ctx, 0, 0)
    # Second pass: build shapes
    _walk(root, ctx, 0, 0)
    return prs


def _collect_rects(parent, ctx, ox, oy):
    """Pre-pass to collect rect positions for text-box matching."""
    for elem in parent:
        tag = _tag(elem)
        if tag == "g":
            tx, ty = _parse_transform(elem.get("transform"))
            _collect_rects(elem, ctx, ox + tx, oy + ty)
        elif tag == "rect":
            tx, ty = _parse_transform(elem.get("transform"))
            x = _num(_ga(elem, "x")) + ox + tx
            y = _num(_ga(elem, "y")) + oy + ty
            w = _num(_ga(elem, "width"))
            h = _num(_ga(elem, "height"))
            if w > 0 and h > 0:
                ctx["rects"].append((x, y, w, h))


def _find_containing_rect(x, y, rects, vb_w, vb_h):
    """Find the smallest rect that contains point (x, y), excluding full-bleed bg."""
    best = None
    best_area = float("inf")
    for rx, ry, rw, rh in rects:
        # Skip full-bleed background rects
        if rw >= vb_w * 0.95 and rh >= vb_h * 0.95:
            continue
        # Skip very small rects (decorative dots, borders)
        if rw < 10 or rh < 10:
            continue
        # Check if text origin is inside this rect (with generous margin)
        margin = 20
        if rx - margin <= x <= rx + rw + margin and ry - margin <= y <= ry + rh + margin:
            area = rw * rh
            if area < best_area:
                best = (rx, ry, rw, rh)
                best_area = area
    return best


def _walk(parent, ctx, ox, oy):
    for elem in parent:
        tag = _tag(elem)

        if tag in ("defs", "metadata", "title", "desc", "style",
                   "marker", "clipPath", "mask", "filter",
                   "linearGradient", "radialGradient", "pattern", "symbol"):
            continue

        # Handle transform on any element
        tx, ty = _parse_transform(elem.get("transform"))
        ex, ey = ox + tx, oy + ty

        if tag == "g":
            _walk(elem, ctx, ex, ey)
            continue

        # Skip invisible elements
        if _opacity(elem) < 0.01:
            continue
        if _ga(elem, "display") == "none" or _ga(elem, "visibility") == "hidden":
            continue

        try:
            if tag == "rect":
                _do_rect(elem, ctx, ex, ey)
            elif tag == "text":
                _do_text(elem, ctx, ex, ey)
            elif tag == "line":
                _do_line(elem, ctx, ex, ey)
            elif tag == "circle":
                _do_circle(elem, ctx, ex, ey)
            elif tag == "ellipse":
                _do_ellipse(elem, ctx, ex, ey)
            elif tag == "polygon":
                _do_polygon(elem, ctx, ex, ey)
            elif tag == "polyline":
                _do_polyline(elem, ctx, ex, ey)
            elif tag == "path":
                _do_path(elem, ctx, ex, ey)
        except Exception as e:
            print(f"  [pptx] skip <{tag}>: {e}")


# ── Shape handlers ────────────────────────────────────────────────────────────

def _apply_fill(shape, fill_val):
    fc = _color(fill_val)
    if fc and fill_val not in ("none", "transparent"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = fc
    else:
        shape.fill.background()


def _apply_stroke(shape, elem):
    stroke_val = _ga(elem, "stroke")
    sc = _color(stroke_val)
    if sc:
        shape.line.color.rgb = sc
        shape.line.width = Pt(max(0.5, _num(_ga(elem, "stroke-width"), 1) * 0.75))
    else:
        shape.line.fill.background()


def _do_rect(elem, ctx, ox, oy):
    sx, sy, slide = ctx["sx"], ctx["sy"], ctx["slide"]
    x = _num(_ga(elem, "x")) + ox
    y = _num(_ga(elem, "y")) + oy
    w = _num(_ga(elem, "width"))
    h = _num(_ga(elem, "height"))
    if w <= 0 or h <= 0:
        return

    # Skip the first full-bleed background rect -> set slide bg instead
    if not ctx["bg_done"] and w >= ctx["vb_w"] * 0.95 and h >= ctx["vb_h"] * 0.95:
        ctx["bg_done"] = True
        fc = _color(_ga(elem, "fill"))
        if fc:
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = fc
        return

    rx = max(_num(_ga(elem, "rx")), _num(_ga(elem, "ry")))
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x * sx), Inches(y * sy),
        Inches(w * sx), Inches(h * sy),
    )
    _apply_fill(shape, _ga(elem, "fill", "#FFFFFF"))
    _apply_stroke(shape, elem)

    # Adjust rounded-corner radius
    if rx > 0 and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            ratio = rx / min(w, h)
            shape.adjustments[0] = min(0.5, ratio)
        except Exception:
            pass


def _do_text(elem, ctx, ox, oy):
    sx, sy, slide = ctx["sx"], ctx["sy"], ctx["slide"]
    full_text = "".join(elem.itertext()).strip()
    if not full_text:
        return

    # ── Gather tspan data (with per-tspan positioning) ──
    tspans = [ts for ts in elem if _tag(ts) == "tspan"]

    lines_data = []
    if tspans:
        for ts in tspans:
            t = "".join(ts.itertext()).strip()
            if not t:
                continue
            ts_x = ts.get("x")
            ts_y = ts.get("y")
            ts_dy = ts.get("dy")
            lines_data.append({
                "text": t,
                "x": _num(ts_x) if ts_x else None,
                "y": _num(ts_y) if ts_y else None,
                "dy": _num(ts_dy) if ts_dy else None,
                "font_size": _ga(ts, "font-size", ""),
                "font_weight": _ga(ts, "font-weight", ""),
                "fill": _ga(ts, "fill", ""),
                "font_family": _ga(ts, "font-family", ""),
            })
        if not lines_data:
            lines_data = [{"text": full_text}]
    else:
        lines_data = [{"text": full_text}]

    lines = [ld["text"] for ld in lines_data]

    # ── Position from <text> element ──
    text_x = _num(_ga(elem, "x")) + ox
    text_y = _num(_ga(elem, "y")) + oy

    if lines_data[0].get("x") is not None:
        text_x = lines_data[0]["x"] + ox
    if lines_data[0].get("y") is not None:
        text_y = lines_data[0]["y"] + oy

    # ── Font properties ──
    font_size = _num(_ga(elem, "font-size"), 20)
    fw = _ga(elem, "font-weight", "")
    bold = fw in ("bold", "700", "800", "900")
    if not bold and fw:
        try:
            bold = int(fw) >= 600
        except Exception:
            pass
    anchor = _ga(elem, "text-anchor", "start")
    fill_val = _ga(elem, "fill") or ctx["def_tc"]
    ff = _ga(elem, "font-family", "Arial")
    ff = ff.strip("'\"").split(",")[0].strip("'\" ")
    if not ff:
        ff = "Arial"

    # Convert font size: SVG units -> PPTX points
    s_avg = (sx + sy) / 2
    pts = max(6, round(font_size * s_avg * 72))

    # ── Try to find containing rect ──
    container = _find_containing_rect(text_x, text_y, ctx["rects"],
                                      ctx["vb_w"], ctx["vb_h"])

    if container:
        rx, ry, rw, rh = container
        xi = rx * sx
        yi = ry * sy
        wi = rw * sx
        hi = rh * sy
        # Internal padding
        pad_x = min(0.12, wi * 0.06)
        pad_y = min(0.08, hi * 0.06)
        xi += pad_x
        yi += pad_y
        wi -= 2 * pad_x
        hi -= 2 * pad_y

        # Clamp font size so text fits inside the container
        # Rough estimate: each line needs ~pts/72 inches height
        available_h_inches = hi - 0.1
        needed_h_inches = len(lines) * (pts / 72) * 1.3
        if needed_h_inches > available_h_inches and available_h_inches > 0:
            scale = available_h_inches / needed_h_inches
            pts = max(6, round(pts * scale))

        # Also check width: estimate chars per line
        max_chars = max(len(l) for l in lines)
        char_w_est = pts / 72 * 0.55  # inches per char at this point size
        needed_w = max_chars * char_w_est
        if needed_w > wi and wi > 0:
            scale = wi / needed_w
            pts = max(6, round(pts * scale))

    else:
        # Estimate text box size from text content
        max_chars = max(len(l) for l in lines)
        char_w_inches = font_size * s_avg * 0.6
        est_w = max(1.0, min(SLIDE_W - 0.4, max_chars * char_w_inches + 0.6))
        line_h_inches = font_size * s_avg * 1.4
        est_h = max(0.4, len(lines) * line_h_inches + 0.2)

        # SVG y = baseline → offset up
        xi = text_x * sx
        yi = max(0, text_y * sy - font_size * sy * 1.0)

        if anchor == "middle":
            xi = max(0, xi - est_w / 2)
        elif anchor == "end":
            xi = max(0, xi - est_w)

        # Clamp within slide
        xi = max(0, min(xi, SLIDE_W - 0.5))
        yi = max(0, min(yi, SLIDE_H - 0.3))
        wi = max(0.3, min(est_w, SLIDE_W - xi))
        hi = max(0.3, min(est_h, SLIDE_H - yi))

    # Clamp final dimensions
    wi = max(0.2, min(wi, SLIDE_W))
    hi = max(0.2, min(hi, SLIDE_H))

    # Create text box
    txBox = slide.shapes.add_textbox(Inches(xi), Inches(yi), Inches(wi), Inches(hi))
    tf = txBox.text_frame
    tf.word_wrap = True
    # Use TEXT_TO_FIT_SHAPE so PowerPoint auto-shrinks if text overflows
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    # Vertical alignment: center text in box
    try:
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")
    except Exception:
        pass

    fc = _color(fill_val) or RGBColor(0x1A, 0x1A, 0x2E)
    align = PP_ALIGN.CENTER if anchor == "middle" else PP_ALIGN.RIGHT if anchor == "end" else PP_ALIGN.LEFT

    for i, ld in enumerate(lines_data):
        line_text = ld["text"]
        if not line_text:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        # Per-line font overrides
        line_pts = pts
        line_bold = bold
        line_fc = fc
        line_ff = ff

        if ld.get("font_size"):
            line_pts = max(6, round(_num(ld["font_size"]) * s_avg * 72))
            # Also clamp overridden font size for container
            if container:
                line_pts = min(line_pts, pts)
        if ld.get("font_weight"):
            lfw = ld["font_weight"]
            line_bold = lfw in ("bold", "700", "800", "900")
            if not line_bold:
                try:
                    line_bold = int(lfw) >= 600
                except Exception:
                    pass
        if ld.get("fill"):
            lc = _color(ld["fill"])
            if lc:
                line_fc = lc
        if ld.get("font_family"):
            lff = ld["font_family"].strip("'\"").split(",")[0].strip("'\" ")
            if lff:
                line_ff = lff

        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(line_pts)
        run.font.name = line_ff
        run.font.bold = line_bold
        run.font.color.rgb = line_fc


def _do_line(elem, ctx, ox, oy):
    sx, sy, slide = ctx["sx"], ctx["sy"], ctx["slide"]
    x1 = (_num(_ga(elem, "x1")) + ox) * sx
    y1 = (_num(_ga(elem, "y1")) + oy) * sy
    x2 = (_num(_ga(elem, "x2")) + ox) * sx
    y2 = (_num(_ga(elem, "y2")) + oy) * sy

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )

    stroke = _ga(elem, "stroke", "#000000")
    sc = _color(stroke)
    if sc:
        connector.line.color.rgb = sc
    connector.line.width = Pt(max(0.5, _num(_ga(elem, "stroke-width"), 1) * 0.75))

    # Dash pattern
    dash = _ga(elem, "stroke-dasharray")
    if dash and dash != "none":
        try:
            ln = connector._element.find(qn("p:spPr")).find(qn("a:ln"))
            if ln is not None:
                prstDash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
                ln.append(prstDash)
        except Exception:
            pass

    # Arrowhead
    marker_end = _ga(elem, "marker-end") or ""
    marker_start = _ga(elem, "marker-start") or ""
    if "url" in marker_end or "marker-end" in elem.get("style", ""):
        _add_arrow(connector, "tail")
    if "url" in marker_start or "marker-start" in elem.get("style", ""):
        _add_arrow(connector, "head")


def _add_arrow(shape, end="tail"):
    """Add arrowhead via OOXML."""
    try:
        spPr = shape._element.find(qn("p:spPr"))
        if spPr is None:
            return
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = spPr.makeelement(qn("a:ln"), {})
            spPr.append(ln)
        tag_name = "a:tailEnd" if end == "tail" else "a:headEnd"
        arr = ln.makeelement(qn(tag_name), {"type": "arrow", "w": "med", "len": "med"})
        ln.append(arr)
    except Exception:
        pass


def _do_circle(elem, ctx, ox, oy):
    sx, sy, slide = ctx["sx"], ctx["sy"], ctx["slide"]
    cx = _num(_ga(elem, "cx")) + ox
    cy = _num(_ga(elem, "cy")) + oy
    r = _num(_ga(elem, "r"))
    if r <= 0:
        return
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches((cx - r) * sx), Inches((cy - r) * sy),
        Inches(2 * r * sx), Inches(2 * r * sy),
    )
    _apply_fill(shape, _ga(elem, "fill", "#000000"))
    _apply_stroke(shape, elem)


def _do_ellipse(elem, ctx, ox, oy):
    sx, sy, slide = ctx["sx"], ctx["sy"], ctx["slide"]
    cx = _num(_ga(elem, "cx")) + ox
    cy = _num(_ga(elem, "cy")) + oy
    rx = _num(_ga(elem, "rx"))
    ry = _num(_ga(elem, "ry"))
    if rx <= 0 or ry <= 0:
        return
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches((cx - rx) * sx), Inches((cy - ry) * sy),
        Inches(2 * rx * sx), Inches(2 * ry * sy),
    )
    _apply_fill(shape, _ga(elem, "fill", "#000000"))
    _apply_stroke(shape, elem)


def _parse_points(pts_str):
    """Parse SVG points attribute -> list of (x,y) floats."""
    nums = [float(n) for n in re.findall(r"[-\d.]+", pts_str)]
    return [(nums[i], nums[i + 1]) for i in range(0, len(nums) - 1, 2)]


def _do_polygon(elem, ctx, ox, oy):
    sx, sy, slide = ctx["sx"], ctx["sy"], ctx["slide"]
    pts = _parse_points(_ga(elem, "points", ""))
    if len(pts) < 3:
        return
    pts = [(x + ox, y + oy) for x, y in pts]
    inch_pts = [(x * sx, y * sy) for x, y in pts]

    builder = slide.shapes.build_freeform(inch_pts[0][0], inch_pts[0][1], scale=EMU_PER_IN)
    builder.add_line_segments(inch_pts[1:])
    builder.add_line_segments([inch_pts[0]])  # close
    shape = builder.convert_to_shape()
    _apply_fill(shape, _ga(elem, "fill", "#000000"))
    _apply_stroke(shape, elem)


def _do_polyline(elem, ctx, ox, oy):
    sx, sy, slide = ctx["sx"], ctx["sy"], ctx["slide"]
    pts = _parse_points(_ga(elem, "points", ""))
    if len(pts) < 2:
        return
    pts = [(x + ox, y + oy) for x, y in pts]
    inch_pts = [(x * sx, y * sy) for x, y in pts]

    builder = slide.shapes.build_freeform(inch_pts[0][0], inch_pts[0][1], scale=EMU_PER_IN)
    builder.add_line_segments(inch_pts[1:])
    shape = builder.convert_to_shape()
    shape.fill.background()
    stroke_val = _ga(elem, "stroke", "#000000")
    sc = _color(stroke_val)
    if sc:
        shape.line.color.rgb = sc
        shape.line.width = Pt(max(0.5, _num(_ga(elem, "stroke-width"), 1) * 0.75))


def _do_path(elem, ctx, ox, oy):
    d = _ga(elem, "d", "")
    if not d:
        return

    # Skip paths with complex curves (they produce garbled output in PPTX)
    # Only allow simple paths with M, L, H, V, Z commands
    has_curves = bool(re.search(r"[CSQTAcsqta]", d))

    sx, sy, slide = ctx["sx"], ctx["sy"], ctx["slide"]

    if has_curves:
        # For curved paths, try to detect if it's a simple shape
        # (e.g., rounded rect) and approximate with more segments
        pts = _parse_path_full(d, ox, oy)
        if len(pts) < 2:
            return
        # Filter out degenerate paths
        xs = [p[0] for p in pts]
        ys = [p[1] for p in pts]
        if max(xs) - min(xs) < 2 and max(ys) - min(ys) < 2:
            return
        # Deduplicate nearly-identical consecutive points to reduce noise
        clean_pts = [pts[0]]
        for p in pts[1:]:
            if abs(p[0] - clean_pts[-1][0]) > 0.5 or abs(p[1] - clean_pts[-1][1]) > 0.5:
                clean_pts.append(p)
        pts = clean_pts
        if len(pts) < 2:
            return
    else:
        pts = _parse_path_full(d, ox, oy)
        if len(pts) < 2:
            return
        xs = [p[0] for p in pts]
        ys = [p[1] for p in pts]
        if max(xs) - min(xs) < 1 and max(ys) - min(ys) < 1:
            return

    inch_pts = [(x * sx, y * sy) for x, y in pts]
    builder = slide.shapes.build_freeform(inch_pts[0][0], inch_pts[0][1], scale=EMU_PER_IN)
    builder.add_line_segments(inch_pts[1:])

    # Close path if d ends with Z
    d_stripped = d.strip()
    if d_stripped and d_stripped[-1] in "Zz":
        builder.add_line_segments([inch_pts[0]])

    shape = builder.convert_to_shape()
    _apply_fill(shape, _ga(elem, "fill", "none"))
    _apply_stroke(shape, elem)


# ── Routes ────────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/generate", methods=["POST"])
def generate():
    if not GEMINI_API_KEY:
        return jsonify({"error": "GEMINI_API_KEY not configured. Set it in environment variables."}), 500
    if not getattr(genai, '_configured', True):
        genai.configure(api_key=GEMINI_API_KEY)

    data = request.json
    text_val = data.get("text", "").strip()
    mode = data.get("mode", "text")
    diag_type = data.get("diagType", "Auto (best fit)")
    theme_bg = data.get("bg", "#FFFFFF")
    theme_text = data.get("text_color", "#1A1A2E")
    accent = data.get("accent", "#0033A0")
    font = data.get("font", "Helvetica Neue")
    image_b64 = data.get("image", None)

    type_note = "" if diag_type == "Auto (best fit)" else f"\nPREFERRED DIAGRAM TYPE: {diag_type}\n"

    prompt = f"""You are a top-tier McKinsey/BCG strategy consultant creating a presentation slide.
Create a clean, professional SVG diagram in the style of McKinsey and BCG consulting decks.

TOPIC: {text_val or '(see image)'}
{type_note}
DESIGN PRINCIPLES (McKinsey/BCG style):
- Action-oriented title: state the insight/conclusion, not just the topic
- Clean grid-based layout with clear visual hierarchy
- Generous white space, no clutter
- Each box should contain concise text (max 3-4 lines per box)
- Use numbered labels or icons to guide the reader's eye
- Source/footnote at bottom-left in small gray text

=== STRICT SVG RULES (MUST FOLLOW) ===

OUTPUT: Only the raw <svg> element. No markdown, no explanation, no code fences.

CANVAS:
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1920 1080">
Do NOT add width or height attributes on the <svg> tag.
First child: <rect x="0" y="0" width="1920" height="1080" fill="{theme_bg}"/>

ALLOWED ELEMENTS ONLY (DO NOT use any other SVG elements):
- <rect> for boxes (use rx="8" for rounded corners)
- <text> with <tspan> for all text content
- <line> for straight connectors and arrows
- <circle> for bullet points or decorative dots
- <ellipse> for ovals
- <polygon> for arrowheads and triangular shapes

FORBIDDEN ELEMENTS (NEVER use these):
- <path> — ABSOLUTELY FORBIDDEN. Never use path elements. They cause rendering errors.
- <polyline> — use <line> instead
- <use>, <image>, <foreignObject> — forbidden
- CSS classes, <style> blocks, gradients, filters, clipPath, mask — forbidden

COLORS:
- Background: {theme_bg}
- Primary text: {theme_text}
- Accent: {accent} (for key boxes, headers, highlights — use sparingly)
- Use lighter tints of the accent for secondary box fills (add opacity="0.1" or "0.15" on accent-colored rects)
- Gray palette for borders/dividers: #E2E8F0, #CBD5E1, #94A3B8

TYPOGRAPHY:
- font-family="'{font}', Arial, sans-serif" on every <text> element
- Title: font-size="44" to "52", font-weight="bold"
- Subtitle/section headers: font-size="28" to "32", font-weight="bold"
- Body text in boxes: font-size="18" to "22"
- Labels/captions: font-size="14" to "16"
- Footer/source: font-size="14", fill="#94A3B8"

TEXT PLACEMENT RULES:
- Every <text> MUST have text-anchor="middle" when centered in a box
- For text inside a <rect>, the text x must equal rect_x + rect_width/2
- For text inside a <rect>, the text y must be vertically centered within the rect
- Use <tspan> with dy="24" or dy="28" for multi-line text within a box
- Keep text SHORT: max 6-8 words per line, max 3-4 lines per box
- ALL text content must be in English

LAYOUT:
- 80px padding from all edges
- Title at top: y="70"
- Main content area: y=130 to y=980
- Footer/source at bottom: y="1050"
- Minimum 20px gap between boxes
- Align boxes to a grid (equal widths, equal heights in each row)

ARROWS/CONNECTORS:
- Use <line> elements with stroke and stroke-width for connectors
- For arrowheads, place a small <polygon> triangle at the end of each line
- Example arrow pointing right: <line x1="100" y1="200" x2="250" y2="200" stroke="#94A3B8" stroke-width="2"/>
  followed by <polygon points="250,192 270,200 250,208" fill="#94A3B8"/>

CRITICAL QUALITY CHECKS:
- No text should overflow outside its containing box
- No overlapping text elements
- Every box that contains text must be large enough for the text
- No random decorative elements — every shape must serve a purpose
- The diagram must be immediately readable and professional"""

    try:
        model = genai.GenerativeModel("gemini-2.5-flash")

        if mode == "text" or (mode == "both" and not image_b64):
            result = model.generate_content(prompt)
        else:
            if not image_b64:
                return jsonify({"error": "No image provided"}), 400

            img_data = base64.b64decode(image_b64.split(",")[-1] if "," in image_b64 else image_b64)
            mime = "image/jpeg"
            if image_b64.startswith("data:"):
                mime = image_b64.split(";")[0].split(":")[1]

            img_suffix = (
                f"Convert this image into a professional SVG diagram.\n\n{prompt}"
                if mode == "image"
                else f"Use the image as structural reference. Description: {text_val}\n\n{prompt}"
            )

            img_part = {"mime_type": mime, "data": img_data}
            result = model.generate_content([img_part, img_suffix])

        raw_text = result.text
        svg_str = extract_svg(raw_text)

        if not svg_str:
            return jsonify({"error": "Could not extract SVG from AI response", "raw": raw_text[:500]}), 500

        return jsonify({"svg": svg_str})

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/export/svg", methods=["POST"])
def export_svg():
    svg_str = request.json.get("svg", "")
    if not svg_str:
        return jsonify({"error": "No SVG"}), 400
    buf = io.BytesIO(svg_str.encode("utf-8"))
    buf.seek(0)
    return send_file(buf, mimetype="image/svg+xml", as_attachment=True,
                     download_name=f"diagram-{int(time.time())}.svg")


@app.route("/api/export/pptx", methods=["POST"])
def export_pptx():
    data = request.json
    svg_str = data.get("svg", "")
    bg_color = data.get("bg", "#FFFFFF")
    text_color = data.get("text_color", "#1A1A2E")

    if not svg_str:
        return jsonify({"error": "No SVG provided"}), 400

    try:
        prs = build_editable_pptx(svg_str, bg_color, text_color)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return send_file(
            buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=f"diagram-{int(time.time())}.pptx",
        )
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"PPTX export failed: {e}"}), 500


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5555, debug=True)
